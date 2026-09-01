"""
Reporting Generator
--------------------
Pulls QA scores, mailbox stats, and task status from the DB and produces:
  - A daily/weekly Excel summary (openpyxl)
  - A simple PPT summary deck (python-pptx) for stakeholder reviews

This replaces the manual "prepare reports using Excel and PowerPoint" line
in the JD.
"""
from datetime import datetime, timedelta, timezone
from sqlalchemy.orm import Session
from sqlalchemy import func
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.utils.dataframe import dataframe_to_rows

from db.models import Call, QAScore, MailboxItem, Task, Agent


def _agent_summary_df(db: Session) -> pd.DataFrame:
    from sqlalchemy import case

    rows = (
        db.query(
            Agent.name,
            func.count(QAScore.id).label("calls_scored"),
            func.avg(QAScore.overall_score).label("avg_score"),
            func.sum(case((QAScore.flagged.is_(True), 1), else_=0)).label("flagged_count"),
        )
        .join(Call, Call.agent_id == Agent.id)
        .join(QAScore, QAScore.call_id == Call.id)
        .group_by(Agent.name)
        .all()
    )
    return pd.DataFrame(rows, columns=["Agent", "Calls Scored", "Avg Score", "Flagged Calls"])


def _mailbox_summary_df(db: Session) -> pd.DataFrame:
    rows = (
        db.query(
            MailboxItem.category,
            func.count(MailboxItem.id).label("count"),
        )
        .group_by(MailboxItem.category)
        .all()
    )
    return pd.DataFrame(rows, columns=["Category", "Count"])


def _task_summary_df(db: Session) -> pd.DataFrame:
    rows = (
        db.query(Task.status, func.count(Task.id).label("count"))
        .group_by(Task.status)
        .all()
    )
    return pd.DataFrame(rows, columns=["Status", "Count"])


def generate_excel_report(db: Session, output_path: str = "ops_report.xlsx") -> str:
    wb = Workbook()

    ws1 = wb.active
    ws1.title = "Agent QA Summary"
    _write_df(ws1, _agent_summary_df(db))

    ws2 = wb.create_sheet("Mailbox Summary")
    _write_df(ws2, _mailbox_summary_df(db))

    ws3 = wb.create_sheet("Task Status")
    _write_df(ws3, _task_summary_df(db))

    wb.save(output_path)
    return output_path


def _write_df(ws, df: pd.DataFrame):
    if df.empty:
        ws.append(["No data available"])
        return
    for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True)):
        ws.append(row)
        if r_idx == 0:
            for cell in ws[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    for col in ws.columns:
        max_len = max(len(str(c.value)) for c in col if c.value is not None) if col else 10
        ws.column_dimensions[col[0].column_letter].width = max_len + 4


def generate_pptx_report(db: Session, output_path: str = "ops_report.pptx") -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Delivery Operations - Weekly Report"
    slide.placeholders[1].text = datetime.now(timezone.utc).strftime("%B %d, %Y")

    bullet_layout = prs.slide_layouts[1]

    def add_table_slide(title, df: pd.DataFrame):
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = title
        body = s.placeholders[1]
        tf = body.text_frame
        tf.clear()
        if df.empty:
            tf.text = "No data available"
            return
        for i, row in df.iterrows():
            line = " | ".join(f"{col}: {row[col]}" for col in df.columns)
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line

    add_table_slide("Agent QA Performance", _agent_summary_df(db))
    add_table_slide("Mailbox Query Breakdown", _mailbox_summary_df(db))
    add_table_slide("Task Status Overview", _task_summary_df(db))

    prs.save(output_path)
    return output_path
